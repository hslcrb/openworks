from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from PyPDF2 import PdfMerger
from pathlib import Path

class DocGen:
    @staticmethod
    def create_docx(path: str | Path, title: str, paragraphs: list[str]):
        doc = Document()
        doc.add_heading(title, level=1)
        for p in paragraphs:
            doc.add_paragraph(p)
        p = Path(path)
        p.parent.mkdir(parents=True, exist_ok=True)
        doc.save(str(p))
        return p

    @staticmethod
    def create_xlsx(path: str | Path, sheet_name: str, rows: list[list]):
        wb = Workbook()
        ws = wb.active
        ws.title = sheet_name
        for r in rows:
            ws.append(r)
        p = Path(path)
        p.parent.mkdir(parents=True, exist_ok=True)
        wb.save(str(p))
        return p

    @staticmethod
    def create_pptx(path: str | Path, title: str, bullets: list[str]):
        prs = Presentation()
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        body = slide.shapes.placeholders[1].text_frame
        for b in bullets:
            body.add_paragraph().text = b
        p = Path(path)
        p.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(p))
        return p

    @staticmethod
    def merge_pdfs(paths: list[str], out_path: str | Path):
        merger = PdfMerger()
        for p in paths:
            merger.append(str(p))
        p = Path(out_path)
        p.parent.mkdir(parents=True, exist_ok=True)
        merger.write(str(p))
        merger.close()
        return p
